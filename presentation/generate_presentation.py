"""Generate K8S Monitor office presentation (PPTX).

Run:  py generate_presentation.py
Output: K8SMonitor_Presentation.pptx (in the same folder)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ---------- Theme ----------
NAVY = RGBColor(0x0B, 0x2E, 0x4F)        # Primary
TEAL = RGBColor(0x00, 0x8C, 0xBA)        # Accent
ORANGE = RGBColor(0xF2, 0x7A, 0x1A)      # Highlight
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)       # Light background
GREEN = RGBColor(0x2E, 0xA8, 0x55)
RED = RGBColor(0xC0, 0x3B, 0x2B)
GREY = RGBColor(0x4A, 0x4A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

# ---------- Helpers ----------

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=16, color=GREY,
                bullet_color=TEAL, line_spacing=1.15):
    """items: list of strings, or list of (text, sub_items)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    first = True
    for item in items:
        if isinstance(item, tuple):
            head, subs = item
        else:
            head, subs = item, []
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet glyph
        r1 = p.add_run()
        r1.text = "\u25CF  "  # solid bullet
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = head
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        for sub in subs:
            ps = tf.add_paragraph()
            ps.line_spacing = line_spacing
            ps.level = 1
            rs1 = ps.add_run()
            rs1.text = "\u2013  "
            rs1.font.name = FONT
            rs1.font.size = Pt(size - 2)
            rs1.font.color.rgb = bullet_color
            rs2 = ps.add_run()
            rs2.text = sub
            rs2.font.name = FONT
            rs2.font.size = Pt(size - 2)
            rs2.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
    add_rect(slide, Inches(0), Inches(0.9), Inches(13.333), Inches(0.08), ORANGE)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.6),
             title, size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=TEAL, anchor=MSO_ANCHOR.TOP)


def add_footer(slide, page_no, total):
    add_rect(slide, Inches(0), Inches(7.25), Inches(13.333), Inches(0.25), LIGHT)
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.3),
             "K8S Monitor  |  AI-Powered Kubernetes Self-Healing",
             size=10, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.5), Inches(0.3),
             f"{page_no} / {total}", size=10, color=GREY,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Build ----------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

TOTAL = 12

# ===== Slide 1: Title =====
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_rect(s, Inches(0), Inches(3.0), Inches(13.333), Inches(0.05), ORANGE)
add_text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.6),
         "K8S MONITOR", size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(1.0),
         "AI-Powered Kubernetes Error Detection & Auto-Fix",
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
         "From Pod Failure  \u2192  Root Cause  \u2192  Pull Request  \u2014  Automatically",
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
         ".NET 8  \u2022  Kubernetes (Kind, local)  \u2022  OpenAI GPT-4o Mini  \u2022  GitHub",
         size=14, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "Internal Technical Overview",
         size=12, color=LIGHT, align=PP_ALIGN.CENTER)

# ===== Slide 2: The Problem =====
s = prs.slides.add_slide(blank)
add_header(s, "The Problem We Solve",
           "Production incidents are slow, repetitive, and expensive")
add_bullets(s, Inches(0.6), Inches(1.7), Inches(6.2), Inches(5.0), [
    ("Pods fail in production, often outside business hours",
     ["Engineers manually tail logs and dig for stack traces"]),
    ("Same bug patterns reappear across services",
     ["Hours lost on root-cause investigation"]),
    ("Fix \u2192 PR \u2192 review cycle is manual and slow",
     ["Mean Time To Resolution (MTTR) is high"]),
    ("Knowledge sits in a few senior engineers' heads",
     ["Hard to scale on-call rotations"]),
], size=17)

# right-side stat card
add_rect(s, Inches(7.4), Inches(1.8), Inches(5.4), Inches(5.0), LIGHT)
add_rect(s, Inches(7.4), Inches(1.8), Inches(5.4), Inches(0.5), TEAL)
add_text(s, Inches(7.4), Inches(1.85), Inches(5.4), Inches(0.4),
         "What \"good\" looks like", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.6), Inches(2.6), Inches(5.0), Inches(0.6),
         "Detect", size=18, bold=True, color=NAVY)
add_text(s, Inches(7.6), Inches(2.95), Inches(5.0), Inches(0.5),
         "errors automatically, 24/7", size=14, color=GREY)
add_text(s, Inches(7.6), Inches(3.7), Inches(5.0), Inches(0.6),
         "Diagnose", size=18, bold=True, color=NAVY)
add_text(s, Inches(7.6), Inches(4.05), Inches(5.0), Inches(0.5),
         "root cause with AI in seconds", size=14, color=GREY)
add_text(s, Inches(7.6), Inches(4.8), Inches(5.0), Inches(0.6),
         "Fix", size=18, bold=True, color=NAVY)
add_text(s, Inches(7.6), Inches(5.15), Inches(5.0), Inches(0.5),
         "open a ready-to-review Pull Request", size=14, color=GREY)
add_footer(s, 2, TOTAL)

# ===== Slide 3: What is K8S Monitor =====
s = prs.slides.add_slide(blank)
add_header(s, "What is K8S Monitor?",
           "A self-healing helper that closes the loop between detection and fix")
add_text(s, Inches(0.6), Inches(1.7), Inches(12.0), Inches(0.7),
         "A .NET 8 job that scans your cluster, asks AI what's wrong, and opens a fix PR \u2014 hands-free.",
         size=20, color=NAVY, bold=True)

# 4 capability cards
cards = [
    ("Scan", "Reads node & pod status, tails pod logs", TEAL),
    ("Detect", "Matches error keywords + stack traces", ORANGE),
    ("Diagnose", "GPT-4o Mini explains root cause", NAVY),
    ("Fix", "Auto-creates a GitHub Pull Request", GREEN),
]
left = 0.6
for i, (h, body, col) in enumerate(cards):
    L = Inches(left + i * 3.1)
    add_rect(s, L, Inches(3.0), Inches(2.9), Inches(2.6), LIGHT)
    add_rect(s, L, Inches(3.0), Inches(2.9), Inches(0.55), col)
    add_text(s, L, Inches(3.02), Inches(2.9), Inches(0.55),
             h, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(left + i * 3.1 + 0.2), Inches(3.75),
             Inches(2.6), Inches(1.8),
             body, size=14, color=GREY)

add_text(s, Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.5),
         "Runs as a Kubernetes CronJob \u2014 every 5 minutes by default.",
         size=14, color=TEAL, bold=True)
add_footer(s, 3, TOTAL)

# ===== Slide 4: How it works =====
s = prs.slides.add_slide(blank)
add_header(s, "How It Works",
           "Six steps, fully automated, on every run")

steps = [
    ("1", "Connect", "Uses kubeconfig to reach the cluster"),
    ("2", "Scan Pods", "Reads tail logs (default 50 lines)"),
    ("3", "Detect", "Errors, exceptions, stack traces"),
    ("4", "Ask AI", "GPT-4o Mini explains + proposes fix"),
    ("5", "Resolve File", "Locates source file in GitHub repo"),
    ("6", "Open PR", "One unified PR per failing pod"),
]
for i, (num, head, body) in enumerate(steps):
    row = i // 3
    col = i % 3
    L = Inches(0.6 + col * 4.15)
    T = Inches(1.9 + row * 2.4)
    add_rect(s, L, T, Inches(4.0), Inches(2.1), LIGHT)
    # circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, L + Inches(0.2), T + Inches(0.3),
                                Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.margin_left = Inches(0); ctf.margin_right = Inches(0)
    ctf.margin_top = Inches(0); ctf.margin_bottom = Inches(0)
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = num
    cr.font.bold = True; cr.font.size = Pt(22); cr.font.color.rgb = WHITE
    cr.font.name = FONT
    add_text(s, L + Inches(1.25), T + Inches(0.25), Inches(2.6), Inches(0.5),
             head, size=18, bold=True, color=NAVY)
    add_text(s, L + Inches(1.25), T + Inches(0.85), Inches(2.6), Inches(1.1),
             body, size=12, color=GREY)
add_footer(s, 4, TOTAL)

# ===== Slide 5: Local Setup with Kind =====
s = prs.slides.add_slide(blank)
add_header(s, "Local Setup \u2014 Kubernetes in Docker (Kind)",
           "We test everything locally on a Kind cluster before any prod use")

add_text(s, Inches(0.6), Inches(1.7), Inches(6.2), Inches(0.5),
         "Why Kind?", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(2.2), Inches(6.2), Inches(4.5), [
    "Spins up a real Kubernetes cluster inside Docker",
    "No cloud cost, no shared environment risk",
    "Same kubeconfig flow as production (EKS / AKS / GKE)",
    "Easy to reset \u2014 delete & recreate in seconds",
    "Perfect for testing failure scenarios safely",
], size=15)

# right card: commands
add_rect(s, Inches(7.2), Inches(1.7), Inches(5.6), Inches(5.0), NAVY)
add_text(s, Inches(7.4), Inches(1.85), Inches(5.2), Inches(0.5),
         "Typical local workflow", size=16, bold=True, color=ORANGE)
cmds = [
    "kind create cluster --name k8s-monitor",
    "kubectl apply -f sample-failing-app.yaml",
    "$env:DRY_RUN = \"true\"   # safe first run",
    "dotnet run",
    "# Review AI output, then DRY_RUN=false",
]
ty = 2.4
for cmd in cmds:
    add_text(s, Inches(7.4), Inches(ty), Inches(5.2), Inches(0.45),
             cmd, size=13, color=LIGHT, font="Consolas")
    ty += 0.45
add_footer(s, 5, TOTAL)

# ===== Slide 6: Pod Log Analysis - What & When =====
s = prs.slides.add_slide(blank)
add_header(s, "Pod Log Analysis",
           "Reading the raw output of a container, directly from the cluster")

add_text(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.5),
         "What it is", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(4.5), [
    "Logs written to stdout/stderr by the running container",
    "Fetched via the Kubernetes API (kubectl logs)",
    "Includes stack traces, startup errors, crashes",
    "Available for every pod \u2014 no app changes needed",
], size=15)

add_text(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.5),
         "When to use it", size=20, bold=True, color=ORANGE)
add_bullets(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.5), [
    "Pod is crashing or in CrashLoopBackOff",
    "Container never starts (config / image error)",
    "Need the most recent failure context, fast",
    "App doesn\u2019t (yet) have proper telemetry",
    "During incidents \u2014 quickest signal available",
], size=15, bullet_color=ORANGE)
add_footer(s, 6, TOTAL)

# ===== Slide 7: Pod Log Analysis - Pros & Cons =====
s = prs.slides.add_slide(blank)
add_header(s, "Pod Log Analysis \u2014 Pros & Cons",
           "Powerful for immediate signal, limited for deep insight")

# Advantages box
add_rect(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.2), LIGHT)
add_rect(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.55), GREEN)
add_text(s, Inches(0.6), Inches(1.72), Inches(6.0), Inches(0.55),
         "  \u2713  Advantages", size=18, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.8), Inches(2.4), Inches(5.6), Inches(4.4), [
    "Zero instrumentation \u2014 works out of the box",
    "Real-time \u2014 see failures as they happen",
    "Captures crashes that telemetry SDKs miss",
    "Stack traces with file + line numbers",
    "Cheap \u2014 no extra ingestion bill",
    "Language- and framework-agnostic",
], size=14, bullet_color=GREEN)

# Disadvantages box
add_rect(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(5.2), LIGHT)
add_rect(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(0.55), RED)
add_text(s, Inches(6.9), Inches(1.72), Inches(6.0), Inches(0.55),
         "  \u2717  Disadvantages", size=18, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.1), Inches(2.4), Inches(5.6), Inches(4.4), [
    "Lost when pod is deleted or restarts (no history)",
    "Unstructured text \u2014 hard to query at scale",
    "No request/user context, no correlation IDs",
    "No business KPIs (latency, success rate)",
    "Noisy \u2014 INFO/DEBUG mixed with real errors",
    "Hard to alert on \u2014 needs external aggregator",
], size=14, bullet_color=RED)
add_footer(s, 7, TOTAL)

# ===== Slide 8: Application Insights =====
s = prs.slides.add_slide(blank)
add_header(s, "Application Insights Logs",
           "Structured, queryable telemetry from inside the application")

add_text(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.5),
         "What it is", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(4.5), [
    "Azure-native telemetry: traces, requests, dependencies",
    "Structured logs with properties, severity, operation IDs",
    "Long retention, KQL queries, dashboards, alerts",
    "End-to-end correlation across microservices",
], size=15)

add_text(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.5),
         "When to use it", size=20, bold=True, color=ORANGE)
add_bullets(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.5), [
    "Performance issues (slow APIs, slow DB calls)",
    "Tracing a request across multiple services",
    "Trend analysis \u2014 \"this started 3 days ago\"",
    "Business metrics \u2014 success rate, P95 latency",
    "Alerting & SLO/SLA reporting",
], size=15, bullet_color=ORANGE)
add_footer(s, 8, TOTAL)

# ===== Slide 9: Comparison Table =====
s = prs.slides.add_slide(blank)
add_header(s, "Pod Logs  vs  Application Insights",
           "They complement each other \u2014 not either/or")

# Build a manual table
cols = ["Aspect", "Pod Log Analysis", "Application Insights"]
rows = [
    ["Setup effort", "None \u2014 built into K8s", "Requires SDK + config"],
    ["Data type", "Raw text", "Structured telemetry"],
    ["Retention", "Minutes \u2013 hours", "Days \u2013 months"],
    ["Best for", "Crashes, startup errors", "Performance, trends, tracing"],
    ["Cost", "Free", "Pay per GB ingested"],
    ["Cross-service view", "No", "Yes (correlation IDs)"],
    ["Available when pod dies", "Limited", "Yes (already shipped)"],
]
top = Inches(1.8)
left = Inches(0.6)
total_w = Inches(12.1)
col_w = [Inches(2.3), Inches(4.7), Inches(5.1)]
row_h = Inches(0.55)

# header
x = left
for i, h in enumerate(cols):
    add_rect(s, x, top, col_w[i], row_h, NAVY)
    add_text(s, x, top, col_w[i], row_h, h, size=15, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[i]

# rows
y = top + row_h
for ri, row in enumerate(rows):
    bg = LIGHT if ri % 2 == 0 else WHITE
    x = left
    for i, val in enumerate(row):
        add_rect(s, x, y, col_w[i], row_h, bg)
        color = NAVY if i == 0 else GREY
        bold = (i == 0)
        add_text(s, x + Inches(0.15), y, col_w[i] - Inches(0.3), row_h,
                 val, size=13, color=color, bold=bold,
                 align=PP_ALIGN.LEFT if i != 0 else PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    y += row_h

add_text(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5),
         "Rule of thumb:  Pod logs catch the failure  \u2192  App Insights explains the journey that led to it.",
         size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_footer(s, 9, TOTAL)

# ===== Slide 10: When to use which =====
s = prs.slides.add_slide(blank)
add_header(s, "When to Use Which?",
           "A simple decision guide for the team")

# Left: Pod logs
add_rect(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.2), LIGHT)
add_rect(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.55), TEAL)
add_text(s, Inches(0.6), Inches(1.72), Inches(6.0), Inches(0.55),
         "Reach for POD LOGS when\u2026", size=18, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(0.8), Inches(2.4), Inches(5.6), Inches(4.4), [
    "A pod is crashing right now",
    "App fails to start (config / secrets / image)",
    "You need a stack trace with file + line",
    "Telemetry is missing, broken, or not yet wired up",
    "Fast triage \u2014 first 60 seconds of an incident",
    "This is exactly what K8S Monitor automates",
], size=14, bullet_color=TEAL)

# Right: AppInsights
add_rect(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(5.2), LIGHT)
add_rect(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(0.55), ORANGE)
add_text(s, Inches(6.9), Inches(1.72), Inches(6.0), Inches(0.55),
         "Reach for APP INSIGHTS when\u2026", size=18, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, Inches(7.1), Inches(2.4), Inches(5.6), Inches(4.4), [
    "Investigating slow or intermittent issues",
    "Tracing a request across multiple services",
    "Comparing today vs last week (trends)",
    "Building dashboards / SLO reports",
    "Setting up production alerts",
    "Root cause needs business context (user, tenant)",
], size=14, bullet_color=ORANGE)
add_footer(s, 10, TOTAL)

# ===== Slide 11: Benefits =====
s = prs.slides.add_slide(blank)
add_header(s, "Benefits for the Team",
           "What changes once K8S Monitor is in place")

benefits = [
    ("Faster", "MTTR drops from hours to minutes", TEAL),
    ("Consistent", "Same quality of analysis every time", ORANGE),
    ("Scalable", "Works while the on-call sleeps", GREEN),
    ("Auditable", "Every fix is a reviewable Pull Request", NAVY),
]
for i, (h, body, col) in enumerate(benefits):
    L = Inches(0.6 + i * 3.1)
    add_rect(s, L, Inches(1.9), Inches(2.95), Inches(3.2), LIGHT)
    add_rect(s, L, Inches(1.9), Inches(2.95), Inches(0.7), col)
    add_text(s, L, Inches(1.92), Inches(2.95), Inches(0.7),
             h, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, L + Inches(0.2), Inches(2.9), Inches(2.55), Inches(2.0),
             body, size=14, color=GREY)

# Bottom callout
add_rect(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.3), NAVY)
add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
         "Engineers stay focused on building features.",
         size=18, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6),
         "K8S Monitor handles the repetitive \"read logs, find file, write fix, open PR\" loop.",
         size=14, color=LIGHT)
add_footer(s, 11, TOTAL)

# ===== Slide 12: Closing =====
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_rect(s, Inches(0), Inches(3.4), Inches(13.333), Inches(0.05), ORANGE)
add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.8),
         "Thank You", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.6),
         "Questions, demo, or a walk-through of the code?",
         size=20, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
         "Repo:  K8SMonitor",
         size=16, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
         "Stack:  .NET 8  \u2022  Kind (local K8s)  \u2022  OpenAI GPT-4o Mini  \u2022  Octokit / GitHub",
         size=14, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "Detect  \u2022  Diagnose  \u2022  Fix  \u2014  Automatically",
         size=14, color=ORANGE, align=PP_ALIGN.CENTER, bold=True)

# ---------- Save ----------
out = os.path.join(os.path.dirname(__file__), "K8SMonitor_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
